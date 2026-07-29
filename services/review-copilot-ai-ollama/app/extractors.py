from __future__ import annotations

import csv
import io
import json
from pathlib import Path

from .chunking import TextSegment


class UnsupportedDocumentError(ValueError):
    pass


def _decode(data: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-8", "cp949"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def _pdf(data: bytes) -> list[TextSegment]:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    return [
        TextSegment(text=page.extract_text() or "", locator=f"page {index}")
        for index, page in enumerate(reader.pages, start=1)
    ]


def _docx(data: bytes) -> list[TextSegment]:
    from docx import Document

    document = Document(io.BytesIO(data))
    paragraphs = [p.text.strip() for p in document.paragraphs if p.text.strip()]
    table_text: list[str] = []
    for table in document.tables:
        for row in table.rows:
            values = [cell.text.strip() for cell in row.cells]
            if any(values):
                table_text.append(" | ".join(values))
    text = "\n".join(paragraphs + table_text)
    return [TextSegment(text=text, locator="document")]


def _pptx(data: bytes) -> list[TextSegment]:
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(data))
    segments: list[TextSegment] = []
    for index, slide in enumerate(presentation.slides, start=1):
        lines: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                value = shape.text.strip()
                if value:
                    lines.append(value)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    values = [cell.text.strip() for cell in row.cells]
                    if any(values):
                        lines.append(" | ".join(values))
        segments.append(TextSegment(text="\n".join(lines), locator=f"slide {index}"))
    return segments


def _xlsx(data: bytes) -> list[TextSegment]:
    from openpyxl import load_workbook

    workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    segments: list[TextSegment] = []
    try:
        for worksheet in workbook.worksheets:
            batch: list[str] = []
            batch_start = 1
            last_row = 1
            for row_number, row in enumerate(worksheet.iter_rows(values_only=True), start=1):
                values = ["" if value is None else str(value) for value in row]
                if not any(value.strip() for value in values):
                    continue
                if not batch:
                    batch_start = row_number
                batch.append(" | ".join(values))
                last_row = row_number
                if len(batch) >= 40:
                    segments.append(
                        TextSegment(
                            text="\n".join(batch),
                            locator=f"sheet {worksheet.title}, rows {batch_start}-{last_row}",
                            metadata={"sheet": worksheet.title},
                        )
                    )
                    batch = []
            if batch:
                segments.append(
                    TextSegment(
                        text="\n".join(batch),
                        locator=f"sheet {worksheet.title}, rows {batch_start}-{last_row}",
                        metadata={"sheet": worksheet.title},
                    )
                )
    finally:
        workbook.close()
    return segments


def _csv(data: bytes) -> list[TextSegment]:
    text = _decode(data)
    reader = csv.reader(io.StringIO(text))
    lines = [" | ".join(row) for row in reader]
    return [TextSegment(text="\n".join(lines), locator="rows")]


def _json(data: bytes) -> list[TextSegment]:
    parsed = json.loads(_decode(data))
    return [
        TextSegment(
            text=json.dumps(parsed, ensure_ascii=False, indent=2),
            locator="document",
        )
    ]


def extract_segments(
    filename: str,
    content_type: str | None,
    data: bytes,
) -> list[TextSegment]:
    suffix = Path(filename).suffix.lower()
    if suffix == ".pdf":
        return _pdf(data)
    if suffix == ".docx":
        return _docx(data)
    if suffix == ".pptx":
        return _pptx(data)
    if suffix == ".xlsx":
        return _xlsx(data)
    if suffix == ".csv":
        return _csv(data)
    if suffix == ".json":
        return _json(data)
    if suffix in {".txt", ".md"}:
        return [TextSegment(text=_decode(data), locator="document")]
    raise UnsupportedDocumentError(
        f"Unsupported document type: {suffix or content_type or 'unknown'}. "
        "Supported: pdf, docx, pptx, xlsx, csv, json, txt, md."
    )

